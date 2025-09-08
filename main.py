import os
import openai
import PyPDF2
from supabase import create_client, Client
import json
from dotenv import load_dotenv
import docx
import pandas as pd
from pptx import Presentation
import re
import nltk
from nltk.tokenize import sent_tokenize
from nltk.corpus import stopwords
import string
import tiktoken

load_dotenv()
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
TABLE_NAME = "people_pal_documents"

if not SUPABASE_URL or not SUPABASE_KEY or not OPENAI_API_KEY:
    raise ValueError("Missing SUPABASE_URL, SUPABASE_KEY, or OPENAI_API_KEY in environment variables.")

openai.api_key = OPENAI_API_KEY

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
MIN_CHUNK_SIZE = 100
MAX_TOKENS = 8000
MAX_CHUNKS_PER_FILE = 1000
MAX_FILE_SIZE_MB = 50

try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')

# --- SUPABASE CLIENT ---
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
encoding = tiktoken.get_encoding("cl100k_base")

def extract_text_from_pdf(pdf_path):
    text = os.path.basename(pdf_path) + "\n"
    with open(pdf_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_txt(txt_path):
    with open(txt_path, "r", encoding="utf-8") as f:
        return os.path.basename(txt_path) + "\n" + f.read()

def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    text = os.path.basename(docx_path) + "\n"
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_csv(csv_path):
    import pandas as pd
    
    file_size_mb = os.path.getsize(csv_path) / (1024 * 1024)
    if file_size_mb > MAX_FILE_SIZE_MB:
        print(f"Skipping large file {csv_path} ({file_size_mb:.1f}MB > {MAX_FILE_SIZE_MB}MB limit)")
        return None
    
    df = pd.read_csv(csv_path)
    text = os.path.basename(csv_path) + "\n"
    
    if len(df) > 1000:
        text += f"Large dataset with {len(df)} rows. Sampling first 500 and last 500 rows.\n"
        df_sample = pd.concat([df.head(500), df.tail(500)])
    else:
        df_sample = df
    
    text += df_sample.to_string(index=False, max_rows=None)
    text += "\n"
    
    return text

def extract_text_from_xlsx(xlsx_path):
    import pandas as pd
    
    file_size_mb = os.path.getsize(xlsx_path) / (1024 * 1024)
    if file_size_mb > MAX_FILE_SIZE_MB:
        print(f"Skipping large file {xlsx_path} ({file_size_mb:.1f}MB > {MAX_FILE_SIZE_MB}MB limit)")
        return None
    
    df_dict = pd.read_excel(xlsx_path, sheet_name=None)
    text = os.path.basename(xlsx_path) + "\n"
    
    for sheet_name, df in df_dict.items():
        text += f"--- Sheet: {sheet_name} ---\n"
        
        if len(df) > 1000:
            text += f"Large dataset with {len(df)} rows. Sampling first 500 and last 500 rows.\n"
            df_sample = pd.concat([df.head(500), df.tail(500)])
        else:
            df_sample = df
        
        text += df_sample.to_string(index=False, max_rows=None)
        text += "\n"
    
    return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = os.path.basename(pptx_path) + "\n"
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".txt":
        return extract_text_from_txt(filepath)
    elif ext == ".docx":
        return extract_text_from_docx(filepath)
    elif ext == ".csv":
        return extract_text_from_csv(filepath)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    else:
        print(f"Skipping unsupported file type: {filepath}")
        return None

def count_tokens(text):
    return len(encoding.encode(text))

def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}\"\'\n]', '', text)
    text = text.strip()
    return text

def extract_title_and_keywords(text):
    lines = text.split('\n')
    title = lines[0] if lines else ""
    
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    stop_words = set(stopwords.words('english'))
    keywords = [word for word in words if word not in stop_words and len(word) > 3]
    
    word_freq = {}
    for word in keywords:
        word_freq[word] = word_freq.get(word, 0) + 1
    
    top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
    return title.strip(), [word for word, freq in top_keywords]

def chunk_text_semantic(text, chunk_size=CHUNK_SIZE, overlap_size=CHUNK_OVERLAP):
    cleaned_text = clean_text(text)
    sentences = sent_tokenize(cleaned_text)
    
    if not sentences:
        return []
    
    chunks = []
    current_chunk = ""
    current_size = 0
    
    for i, sentence in enumerate(sentences):
        sentence_length = len(sentence)
        
        if current_size + sentence_length > chunk_size and current_chunk:
            if len(current_chunk.strip()) >= MIN_CHUNK_SIZE:
                chunks.append(current_chunk.strip())
            
            overlap_text = ""
            overlap_size_current = 0
            for j in range(i-1, -1, -1):
                if overlap_size_current + len(sentences[j]) <= overlap_size:
                    overlap_text = sentences[j] + " " + overlap_text
                    overlap_size_current += len(sentences[j])
                else:
                    break
            
            current_chunk = overlap_text + sentence + " "
            current_size = len(current_chunk)
        else:
            current_chunk += sentence + " "
            current_size += sentence_length + 1
    
    if current_chunk.strip() and len(current_chunk.strip()) >= MIN_CHUNK_SIZE:
        chunks.append(current_chunk.strip())
    
    return [chunk for chunk in chunks if chunk.strip()]

def split_text_by_tokens(text, max_tokens=MAX_TOKENS):
    tokens = encoding.encode(text)
    if len(tokens) <= max_tokens:
        return [text]
    
    chunks = []
    current_pos = 0
    
    while current_pos < len(tokens):
        end_pos = min(current_pos + max_tokens, len(tokens))
        chunk_tokens = tokens[current_pos:end_pos]
        chunk_text = encoding.decode(chunk_tokens)
        chunks.append(chunk_text)
        current_pos = end_pos
    
    return chunks

def get_embedding(text):
    cleaned_text = clean_text(text)
    if not cleaned_text or len(cleaned_text.strip()) < 10:
        return None
    
    token_count = count_tokens(cleaned_text)
    if token_count > MAX_TOKENS:
        print(f"Text too long ({token_count} tokens), truncating to {MAX_TOKENS} tokens")
        tokens = encoding.encode(cleaned_text)
        truncated_tokens = tokens[:MAX_TOKENS]
        cleaned_text = encoding.decode(truncated_tokens)
    
    response = openai.embeddings.create(
        input=cleaned_text,
        model="text-embedding-3-small"
    )
    return response.data[0].embedding

def upload_to_supabase(content, embedding, metadata=None):
    data = {
        "content": content,
        "embedding": embedding,
        "metadata": metadata or {}
    }
    supabase.table(TABLE_NAME).insert(data).execute()

def main():
    filedump_dir = "filedump"
    for root, dirs, files in os.walk(filedump_dir):
        for filename in files:
            filepath = os.path.join(root, filename)
            relative_path = os.path.relpath(filepath, filedump_dir)
            text = extract_text_from_file(filepath)
            if not text:
                continue
            try:
                title, keywords = extract_title_and_keywords(text)
                semantic_chunks = chunk_text_semantic(text)
                
                final_chunks = []
                for semantic_chunk in semantic_chunks:
                    if count_tokens(semantic_chunk) > MAX_TOKENS:
                        token_split_chunks = split_text_by_tokens(semantic_chunk)
                        final_chunks.extend(token_split_chunks)
                    else:
                        final_chunks.append(semantic_chunk)
                
                if len(final_chunks) > MAX_CHUNKS_PER_FILE:
                    print(f"Warning: {relative_path} has {len(final_chunks)} chunks, limiting to {MAX_CHUNKS_PER_FILE}")
                    final_chunks = final_chunks[:MAX_CHUNKS_PER_FILE]
                
                total_chunks = len(final_chunks)
                print(f"Processing {relative_path}: {total_chunks} chunks generated (semantic + token-aware)")
                
                for idx, chunk in enumerate(final_chunks):
                    progress_pct = ((idx + 1) / total_chunks) * 100
                    print(f"Embedding and uploading chunk {idx+1}/{total_chunks} ({progress_pct:.1f}%) from {relative_path}")
                    emb = get_embedding(chunk)
                    
                    if emb is None:
                        print(f"Skipping chunk {idx+1} - too short or empty after cleaning")
                        continue
                    
                    chunk_title, chunk_keywords = extract_title_and_keywords(chunk)
                    
                    metadata = {
                        "chunk_index": idx,
                        "total_chunks": total_chunks,
                        "source_file": relative_path,
                        "file_title": title,
                        "file_keywords": keywords,
                        "chunk_title": chunk_title,
                        "chunk_keywords": chunk_keywords,
                        "chunk_length": len(chunk),
                        "file_type": os.path.splitext(filename)[1].lower()
                    }
                    upload_to_supabase(chunk, emb, metadata)
                    
                print(f"Finished processing {relative_path}, deleting file.")
                os.remove(filepath)
            except Exception as e:
                print(f"Error processing {relative_path}: {e}")
    print("All files in filedump processed, uploaded, and deleted.")

if __name__ == "__main__":
    main()
